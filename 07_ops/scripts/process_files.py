#!/usr/bin/env python3
"""
Unified file extractor for project brain.

Extracts text from PDF, DOCX, and PPTX files into markdown format.
Includes OCR fallback for scanned PDFs.

Usage:
    python process_files.py <file_or_directory> [--output <dir>]
"""

import argparse
import logging
import sys
from datetime import date
from pathlib import Path

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(levelname)s: %(message)s"
)
logger = logging.getLogger(__name__)


def extract_pdf(file_path: Path) -> tuple[str, dict]:
    """Extract text from PDF. Falls back to OCR if no text found."""
    from pypdf import PdfReader
    
    metadata = {"pages": 0}
    
    try:
        reader = PdfReader(file_path)
        metadata["pages"] = len(reader.pages)
        
        text_parts = []
        for i, page in enumerate(reader.pages, 1):
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(f"## Page {i}\n\n{page_text}")
        
        text = "\n\n".join(text_parts)
        
        # If no text extracted, try OCR
        if not text.strip():
            logger.info(f"No text found in PDF, attempting OCR: {file_path.name}")
            text, ocr_metadata = extract_pdf_ocr(file_path)
            metadata.update(ocr_metadata)
            metadata["extraction_method"] = "ocr"
        else:
            metadata["extraction_method"] = "text"
        
        return text, metadata
        
    except Exception as e:
        logger.error(f"Failed to extract PDF {file_path.name}: {e}")
        return "", metadata


def extract_pdf_ocr(file_path: Path) -> tuple[str, dict]:
    """Extract text from scanned PDF using OCR."""
    metadata = {}
    
    try:
        from pdf2image import convert_from_path
        import pytesseract
        
        images = convert_from_path(file_path)
        metadata["pages"] = len(images)
        
        text_parts = []
        for i, image in enumerate(images, 1):
            page_text = pytesseract.image_to_string(image)
            if page_text.strip():
                text_parts.append(f"## Page {i}\n\n{page_text}")
        
        return "\n\n".join(text_parts), metadata
        
    except ImportError as e:
        logger.warning(f"OCR dependencies not available: {e}")
        logger.info("Install pdf2image and pytesseract for OCR support")
        return "", metadata
    except Exception as e:
        logger.error(f"OCR failed for {file_path.name}: {e}")
        return "", metadata


def extract_docx(file_path: Path) -> tuple[str, dict]:
    """Extract text from Word document."""
    from docx import Document
    
    metadata = {}
    
    try:
        doc = Document(file_path)
        
        text_parts = []
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                # Preserve heading styles
                if para.style.name.startswith("Heading"):
                    level = para.style.name.replace("Heading ", "")
                    try:
                        hashes = "#" * int(level)
                    except ValueError:
                        hashes = "##"
                    text_parts.append(f"{hashes} {para.text}")
                else:
                    text_parts.append(para.text)
        
        # Extract tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                table_text.append(f"| {row_text} |")
            if table_text:
                # Add header separator after first row
                if len(table_text) > 1:
                    cols = len(table.rows[0].cells)
                    separator = "| " + " | ".join(["---"] * cols) + " |"
                    table_text.insert(1, separator)
                text_parts.append("\n".join(table_text))
        
        metadata["paragraphs"] = len(doc.paragraphs)
        metadata["tables"] = len(doc.tables)
        
        return "\n\n".join(text_parts), metadata
        
    except Exception as e:
        logger.error(f"Failed to extract DOCX {file_path.name}: {e}")
        return "", metadata


def extract_pptx(file_path: Path) -> tuple[str, dict]:
    """Extract text from PowerPoint presentation."""
    from pptx import Presentation
    
    metadata = {}
    
    try:
        prs = Presentation(file_path)
        metadata["slides"] = len(prs.slides)
        
        text_parts = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = [f"## Slide {i}"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Extract table content
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells
                        )
                        slide_text.append(f"| {row_text} |")
            
            # Extract notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    slide_text.append(f"\n**Notes:** {notes}")
            
            if len(slide_text) > 1:  # More than just the header
                text_parts.append("\n\n".join(slide_text))
        
        return "\n\n---\n\n".join(text_parts), metadata
        
    except Exception as e:
        logger.error(f"Failed to extract PPTX {file_path.name}: {e}")
        return "", metadata


def process_file(file_path: Path, output_dir: Path | None = None) -> bool:
    """Process a single file and save extracted text as markdown."""
    
    # Determine extractor based on extension
    extractors = {
        ".pdf": extract_pdf,
        ".docx": extract_docx,
        ".pptx": extract_pptx,
    }
    
    ext = file_path.suffix.lower()
    if ext not in extractors:
        logger.warning(f"Unsupported file type: {file_path.name}")
        return False
    
    logger.info(f"Processing: {file_path.name}")
    
    # Extract content
    text, metadata = extractors[ext](file_path)
    
    if not text.strip():
        logger.warning(f"No content extracted from: {file_path.name}")
        return False
    
    # Build output path
    if output_dir:
        output_path = output_dir / f"{file_path.stem}.md"
    else:
        output_path = file_path.with_suffix(".md")
    
    # Build frontmatter
    frontmatter_lines = [
        "---",
        f"source_file: {file_path.name}",
        f"extracted: {date.today().isoformat()}",
    ]
    for key, value in metadata.items():
        frontmatter_lines.append(f"{key}: {value}")
    frontmatter_lines.append("---")
    
    # Combine frontmatter and content
    output_content = "\n".join(frontmatter_lines) + f"\n\n# {file_path.stem}\n\n{text}\n"
    
    # Write output
    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(output_content)
        logger.info(f"Created: {output_path}")
        return True
    except Exception as e:
        logger.error(f"Failed to write {output_path}: {e}")
        return False


def main():
    parser = argparse.ArgumentParser(
        description="Extract text from documents (PDF, DOCX, PPTX) into markdown."
    )
    parser.add_argument(
        "path",
        type=Path,
        help="File or directory to process"
    )
    parser.add_argument(
        "--output", "-o",
        type=Path,
        default=None,
        help="Output directory (default: same as source)"
    )
    
    args = parser.parse_args()
    
    if not args.path.exists():
        logger.error(f"Path not found: {args.path}")
        sys.exit(1)
    
    # Collect files to process
    if args.path.is_file():
        files = [args.path]
    else:
        files = list(args.path.glob("*.pdf")) + \
                list(args.path.glob("*.docx")) + \
                list(args.path.glob("*.pptx"))
    
    if not files:
        logger.warning("No supported files found")
        sys.exit(0)
    
    # Process files
    success = 0
    for file_path in files:
        if process_file(file_path, args.output):
            success += 1
    
    logger.info(f"Processed {success}/{len(files)} files")
    
    if success < len(files):
        sys.exit(1)


if __name__ == "__main__":
    main()

